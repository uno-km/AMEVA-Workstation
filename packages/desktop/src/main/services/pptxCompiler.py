# -*- coding: utf-8 -*-
"""
@file pptxCompiler.py
@role PPTX slide compiler service (exports PPTX to PNG images or text outlines)
"""
import sys
import os
import json

def process_pptx(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        try:
            os.makedirs(output_dir)
        except Exception as e:
            return {"success": False, "error": f"Failed to create output directory: {str(e)}"}

    # 1) Try PowerPoint COM API (requires MS PowerPoint installed on Windows)
    try:
        import win32com.client
        # Dispatch PowerPoint Application
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        
        # Open presentation: Open(FileName, ReadOnly, Untitled, WithWindow)
        # We set WithWindow to False to keep PPT backgrounded
        presentation = powerpoint.Presentations.Open(pptx_path, True, False, False)
        
        # Export entire presentation as PNG slide images
        # PowerPoint exports files named Slide1.PNG, Slide2.PNG... in output_dir
        presentation.Export(output_dir, "PNG")
        presentation.Close()
        powerpoint.Quit()

        # Gather exported slide images
        files = sorted(
            [f for f in os.listdir(output_dir) if f.lower().endswith('.png')],
            key=lambda x: int(''.join(filter(str.isdigit, x)) or 0)
        )
        slide_images = [os.path.join(output_dir, f).replace('\\', '/') for f in files]
        
        return {
            "success": True,
            "slides": slide_images,
            "fallback": False
        }
    except Exception as com_err:
        # 2) Fallback to text outline parsing if PowerPoint COM is not available
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            slides_text = []
            for idx, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                slides_text.append({
                    "slide_index": idx + 1,
                    "texts": texts
                })
            return {
                "success": True,
                "slides": [],
                "fallback": True,
                "slides_text": slides_text,
                "error": str(com_err)
            }
        except Exception as fallback_err:
            return {
                "success": False,
                "error": f"COM Error: {str(com_err)} | python-pptx Fallback Error: {str(fallback_err)}"
            }

if __name__ == "__main__":
    # Ensure stdout handles unicode encoding safely on Windows console
    try:
        import sys
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.detach(), encoding='utf-8')
        sys.stderr = io.TextIOWrapper(sys.stderr.detach(), encoding='utf-8')
    except Exception:
        pass

    if len(sys.argv) < 3:
        print(json.dumps({"success": False, "error": "Missing parameters. Usage: python pptxCompiler.py <pptx_path> <output_dir>"}))
        sys.exit(1)

    pptx_input = sys.argv[1]
    pptx_output = sys.argv[2]
    
    result = process_pptx(pptx_input, pptx_output)
    print(json.dumps(result, ensure_ascii=False))
